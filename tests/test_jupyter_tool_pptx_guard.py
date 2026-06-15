"""Regression tests for PPTX creation guardrails in execute_code."""

from __future__ import annotations

import pytest

from box_agent.tools.jupyter_tool import MAX_EXECUTE_CODE_CHARS, JupyterSandboxTool


def test_execute_code_schema_exposes_code_size_limit():
    code_schema = JupyterSandboxTool().parameters["properties"]["code"]

    assert code_schema["maxLength"] == MAX_EXECUTE_CODE_CHARS
    assert "split large scripts" in code_schema["description"]


@pytest.mark.asyncio
async def test_execute_code_rejects_oversized_code_before_kernel_start(monkeypatch):
    def fail_if_sandbox_requested(self):
        raise AssertionError("oversized code should be rejected before sandbox startup")

    monkeypatch.setattr(JupyterSandboxTool, "_get_sandbox_env", fail_if_sandbox_requested)
    tool = JupyterSandboxTool()
    code = "x = 1\n" + ("print(x)\n" * MAX_EXECUTE_CODE_CHARS)

    result = await tool.execute(code=code)

    assert result.success is False
    assert result.error is not None
    assert result.error.startswith("EXECUTE_CODE_TOO_LARGE")
    assert "Split the work into multiple execute_code calls" in result.error


def test_execute_code_blocks_bare_python_pptx_new_deck_constructor():
    code = """\
from pptx import Presentation

prs = Presentation()
prs.save("deck.pptx")
"""

    assert JupyterSandboxTool._looks_like_python_pptx_new_deck(code)


def test_execute_code_allows_existing_pptx_inspection():
    code = """\
from pptx import Presentation

prs = Presentation("existing.pptx")
print(len(prs.slides))
"""

    assert not JupyterSandboxTool._looks_like_python_pptx_new_deck(code)


def test_execute_code_allows_non_pptx_data_work():
    code = """\
import pandas as pd

df = pd.read_csv("data.csv")
print(df.describe())
"""

    assert not JupyterSandboxTool._looks_like_python_pptx_new_deck(code)
